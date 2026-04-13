#!/usr/bin/env python3
"""
Generate PowerPoint presentation from midterm presentation content.
Uses python-pptx to create a professional slide deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# Color scheme matching the HTML presentation
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)  # #1a1a2e
GREEN_ACCENT = RGBColor(0x4E, 0xCC, 0xA3)  # #4ecca3
RED_ACCENT = RGBColor(0xFF, 0x6B, 0x6B)  # #ff6b6b
YELLOW_ACCENT = RGBColor(0xF1, 0xC4, 0x0F)  # #f1c40f
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)

def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle, presenter="Robert Katz", context="VTS Tournament - Stacking Sats Challenge"):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = YELLOW_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Context
    ctx_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.6))
    tf = ctx_box.text_frame
    p = tf.paragraphs[0]
    p.text = context
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Presenter
    pres_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = pres_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{presenter} | GT ID: rkatz8@gatech.edu"
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, image_path=None, note=None):
    """Add a content slide with bullets and optional image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    p.alignment = PP_ALIGN.LEFT
    
    # Content area
    if image_path and os.path.exists(image_path):
        # Split layout: bullets on left, image on right
        content_left = Inches(0.5)
        content_width = Inches(4.5)
        
        # Add image
        try:
            slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.3), width=Inches(4.3))
        except:
            pass  # Image failed to load, continue with text only
    else:
        content_left = Inches(0.8)
        content_width = Inches(8.4)
    
    # Bullets
    bullet_box = slide.shapes.add_textbox(content_left, Inches(1.3), content_width, Inches(5.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.level = 0
        p.space_after = Pt(12)
    
    # Footer note
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content, left_title=None, right_title=None):
    """Add a slide with two columns."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    
    # Left column title
    if left_title:
        lt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(0.5))
        tf = lt_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GREEN_ACCENT
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.3), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    
    # Right column title
    if right_title:
        rt_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.5))
        tf = rt_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GREEN_ACCENT
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    
    return slide

def add_image_slide(prs, title, image_path, caption=None):
    """Add a slide with a full-width image."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    
    # Image
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.3), width=Inches(8.4))
        except Exception as e:
            # Add error message
            err_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1))
            tf = err_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Image: {os.path.basename(image_path)}"
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY
    
    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.6))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(18)
        p.font.color.rgb = YELLOW_ACCENT
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_comparison_table_slide(prs, title, headers, rows, highlight_col=None):
    """Add a slide with a comparison table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    
    # Create table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.8), Inches(1.5), Inches(8.4), Inches(0.6 * num_rows)).table
    
    # Style headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN_ACCENT
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(18)
        paragraph.font.bold = True
        paragraph.font.color.rgb = DARK_BG
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Style rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    
    return slide

def add_quote_slide(prs, title, quote, attribution=None):
    """Add a slide with a prominent quote."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Quote box with left border
    quote_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(3))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    quote_box.line.color.rgb = GREEN_ACCENT
    quote_box.line.width = Pt(4)
    
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # Attribution
    if attribution:
        att_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.6))
        tf = att_box.text_frame
        p = tf.paragraphs[0]
        p.text = attribution
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = GREEN_ACCENT
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_pipeline_slide(prs, title):
    """Add the pipeline visualization slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    
    # Pipeline steps
    steps = [
        ("Price Data", "(90 days)", 1.5),
        ("GAF Image Generation", "", 2.3),
        ("Deep CNN", "(296K params)", 3.1),
        ("Temperature Calibration", "", 3.9),
        ("Tilt-Based Allocation", "", 4.7),
        ("EMA Smoothing", "", 5.5),
        ("Normalized Weights", "", 6.3),
    ]
    
    for step, detail, top in steps:
        # Step box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(top), Inches(4), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
        box.line.color.rgb = GREEN_ACCENT
        
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow (except for last step)
        if detail:
            arrow = slide.shapes.add_textbox(Inches(4.5), Inches(top + 0.6), Inches(1), Inches(0.15))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.text = "↓"
            p.font.size = Pt(20)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER
    
    # Hypothesis
    hyp_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.5))
    tf = hyp_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Hypothesis: Visual patterns in GAF images predict returns under strict causality constraints"
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    """Generate the complete PowerPoint presentation."""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    images_dir = "images"
    
    # SLIDE 1: Title Slide
    add_title_slide(prs, 
        "When AI Meets Market Reality",
        "Midterm Finding: Baseline Beats CNN (So Far)",
        presenter="Robert Katz")
    
    # SLIDE 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "Sponsor & Problem: VTS Tournament - Stacking Sats Challenge",
        "   Help retail investors accumulate Bitcoin more effectively than DCA",
        "",
        "Decision This Enables: Daily allocation between BTC and cash",
        "   Optimizing for risk-adjusted returns under strict causality",
        "",
        "Success Definition: Top 60% percentile in recency-weighted SPD",
        "   Beat DCA consistently; demonstrate scientific rigor"
    ])
    
    # SLIDE 3: Why Visual Trading Systems
    add_content_slide(prs, "Why Visual Trading Systems?", [
        "My Finance Foundation (1970s MBA Era)",
        "   • Magee & Edwards' Technical Analysis of Stock Trends",
        "   • Human traders 'see' patterns quants often miss",
        "",
        "The Discovery (OMSA Program)",
        "   • JP Morgan AI group research with Professor Baich",
        "   • CNNs achieving 0.90+ F1 on candlestick patterns",
        "",
        "The Gap & The Test",
        "   • Pattern recognition ≠ Profitable trading",
        "   • Does academic success translate under real causality?"
    ])
    
    # SLIDE 4: Tournament Challenge
    add_content_slide(prs, "The Tournament Challenge", [
        "Objective: Predict daily BTC allocation weights (2016-2025)",
        "",
        "Metric: Recency-Weighted SPD Percentile",
        "   (tournament rank vs all submissions, emphasizing recent periods;",
        "    higher is better; 50th ≈ median submission performance)",
        "",
        "Constraints: Σwᵢ = 1.0, wᵢ ≥ 1e-5",
        "",
        "Critical: Strict causality - no lookahead allowed",
        "",
        "Target: Top performers >60% percentile",
        "",
        "Causality Verification: Last-row test confirmed zero leakage"
    ], note="Evaluation harness + leakage tests automated")
    
    # SLIDE 5: Data Pipeline
    add_content_slide(prs, "Data Pipeline & Methodology", [
        "Dataset: BTC daily OHLCV (2016-2025) from Polygon",
        "",
        "Feature Construction: GAF images from 90-day rolling windows",
        "",
        "Pipeline: Raw Data → GAF Images → CNN (296K) → Allocator → Output",
        "",
        "Training: 2014-2015 (public BTC OHLCV)",
        "Evaluation: Strictly 2016-2025 (tournament period), no overlap",
        "Leakage Controls: Strict temporal embargo + causality tests"
    ])
    
    # SLIDE 5b: Execution Assumptions & Limitations
    add_content_slide(prs, "Execution Assumptions & Limitations", [
        "Signal Timestamp: End of day t (uses data through close t)",
        "Execution Timestamp: Day t (same-day approximation)",
        "Execution Price: Close t (proxy for next-open or VWAP)",
        "",
        "Transaction Costs: 0 bps base case; viable to 50+ bps",
        "Look-ahead Bias: NONE (3 validation tests passed)",
        "Cash Constraints: Budget-normalized (dynamic = naive total)",
        "Sell Logic: NONE (buy-and-hold accumulation only)",
        "",
        "Validation: Last-row modification, purge test, shift test"
    ], note="See EXECUTION_ASSUMPTIONS.md for full sensitivity analysis")
    
    # SLIDE 6: Our Approach (Pipeline)
    add_pipeline_slide(prs, "Our Initial Approach")
    
    # SLIDE 7: First Results - SCOREBOARD
    add_comparison_table_slide(prs, "Reality Check: Scoreboard",
        ["Metric", "CNN", "Simplified", "Target"],
        [
            ["RW SPD Percentile", "41.43%", "41.94%", ">60%"],
            ["Win Rate vs DCA", "54.32%", "70.42%", ">50%"],
            ["Windows Evaluated", "3,076", "3,076", "—"]
        ]
    )
    
    # Add scoreboard summary box
    slide = prs.slides[-1]
    score_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.0), Inches(7), Inches(0.6))
    score_box.fill.solid()
    score_box.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    score_box.line.color.rgb = GREEN_ACCENT
    tf = score_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CNN: 41.43%  |  Simplified: 41.94%  |  Target: >60%"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add win rate paradox explanation
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.5))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Note: Higher win rate can still yield lower SPD because SPD rewards magnitude/timing, not just 'up days'"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER
    
    # SLIDE 8: Progress & Timeline
    add_two_column_slide(prs, "Current Status vs Plan",
        left_content=[
            "DONE (Midterm):",
            "• Working baseline + all tests passing",
            "• GAF image pipeline",
            "• CNN architecture (296K)",
            "• Causality verification",
            "• Ablation harness",
            "• Baseline comparison"
        ],
        right_content=[
            "Phase I (to March 2026):",
            "   • Weekly granularity testing",
            "   • Regime detection model",
            "   • XGBoost baseline",
            "",
            "Phase II (to Final):",
            "   • Ensemble integration",
            "   • Full backtest validation",
            "   • Final documentation"
        ],
        left_title="Complete",
        right_title="Next Steps"
    )
    
    # SLIDE 9: Ablation Study #1
    img_path = os.path.join(images_dir, "01_ablation_sensitivity.png")
    slide = add_image_slide(prs, "Ablation Study #1: Allocator Sizing", img_path,
        caption="Conclusion: Minimal impact (+0.35 pp max). Problem NOT allocator aggressiveness.")
    
    # Add Ablation #0 banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(6.5), Inches(2.5), Inches(0.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = YELLOW_ACCENT
    banner.line.color.rgb = YELLOW_ACCENT
    tf = banner.text_frame
    p = tf.paragraphs[0]
    p.text = "Ablation #0: Is signal > random?"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BG
    p.alignment = PP_ALIGN.CENTER
    
    # Add Takeaway footer
    takeaway = slide.shapes.add_textbox(Inches(3.0), Inches(6.5), Inches(6.5), Inches(0.5))
    tf = takeaway.text_frame
    p = tf.paragraphs[0]
    p.text = "Takeaway: Allocator tuning doesn't move the needle → problem is NOT sizing."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    p.alignment = PP_ALIGN.LEFT
    
    # SLIDE 10: Ablation Study #2 - The "Aha!" Moment
    img_path = os.path.join(images_dir, "02_ablation_signal_quality.png")
    slide = add_image_slide(prs, "Ablation Study #2: The 'Aha!' Moment", img_path,
        caption="Finding: No predictive signal detected under strict causality. CNN output statistically indistinguishable from a neutral baseline.")
    
    # Add Ablation #0 answer
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(6.5), Inches(2.5), Inches(0.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RED_ACCENT
    banner.line.color.rgb = RED_ACCENT
    tf = banner.text_frame
    p = tf.paragraphs[0]
    p.text = "Ablation #0 Answer: NO"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add Takeaway footer
    takeaway = slide.shapes.add_textbox(Inches(3.0), Inches(6.5), Inches(6.5), Inches(0.5))
    tf = takeaway.text_frame
    p = tf.paragraphs[0]
    p.text = "Takeaway: CNN signal ≈ neutral baseline → no edge detected."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.alignment = PP_ALIGN.LEFT
    
    # SLIDE 11: Ablation Study #3
    slide = add_image_slide(prs, "Ablation Study #3: EMA Smoothing", img_path,
        caption="Result: ALL alphas produce IDENTICAL results. Constant signal → tilt always zero → no smoothing effect")
    
    # Add Takeaway footer
    takeaway = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.4))
    tf = takeaway.text_frame
    p = tf.paragraphs[0]
    p.text = "Takeaway: Constant output → smoothing irrelevant."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # SLIDE 12: Root Causes
    img_path = os.path.join(images_dir, "06_root_causes.png")
    add_image_slide(prs, "Midterm Diagnosis: Working Hypotheses", img_path,
        caption="Takeaway: Daily patterns lack predictive power. Regime-aware / longer-horizon approach needed.")
    
    # SLIDE 13: The Pivot
    add_two_column_slide(prs, "The Evidence-Based Pivot",
        left_content=[
            "CNN Approach:",
            "• 41.43% RW Percentile",
            "• 54.32% Win Rate",
            "• 1.1 MB Model",
            "• ~17 min execution",
            "• Complex, no edge detected"
        ],
        right_content=[
            "Simplified Baseline:",
            "• 41.94% RW Percentile",
            "• 70.42% Win Rate",
            "• No Model",
            "• ~2.5 min execution",
            "• Robust, simple, preferred"
        ],
        left_title="CNN (296K params)",
        right_title="Simplified (Random/Neutral)"
    )
    
    # SLIDE 14: Side-by-Side Comparison
    img_path = os.path.join(images_dir, "04_performance_comparison.png")
    add_image_slide(prs, "Side-by-Side Comparison", img_path,
        caption="Simplified (neutral baseline) wins on ALL dimensions")
    
    # SLIDE 15: What We Did Right
    add_two_column_slide(prs, "What We Did Right",
        left_content=[
            "Technical Excellence:",
            "• Strict causality enforcement",
            "• Deterministic (seed=42)",
            "• Comprehensive testing",
            "• Grader-safe implementation"
        ],
        right_content=[
            "Scientific Rigor:",
            "• Hypothesis-driven development",
            "• Systematic ablation studies",
            "• Transparent negative results",
            "• Learning orientation"
        ]
    )
    
    # SLIDE 16: Key Lessons
    add_content_slide(prs, "Key Lessons Learned", [
        "1. Ablation studies are critical",
        "   Without them, we'd never detect the lack of signal",
        "",
        "2. Test 'better than random?' early (Ablation #0)",
        "   Should be first test, not buried in sequence",
        "",
        "3. Complexity ≠ Performance",
        "   296K parameters matched by neutral baseline",
        "",
        "4. Simple often beats complex",
        "   Occam's Razor applies to trading"
    ])
    
    # SLIDE 16b: Post-Midterm Plan (Evidence-Driven)
    add_content_slide(prs, "Post-Midterm Plan (Evidence-Driven)", [
        "1. Regime-Gated Allocator (Highest Leverage)",
        "   • Choose between DCA / buy-more / buy-less based on causal regime",
        "   • Compare: volatility regime, trend regime, 2-state HMM",
        "   • Action space tiny: 3 policies beat daily-noise problem",
        "",
        "2. Reduce Daily Noise (Signal Horizon)",
        "   • Shift from 1-day direction to 5–10 day horizon signals",
        "   • Smooth/hold constant for weekly rebalancing",
        "",
        "3. Robustness & Validation (Keep it Credible)",
        "   • Walk-forward evaluation with purged splits",
        "   • Report performance by subperiod + regime",
        "",
        "Stretch: Point-and-Figure → simple features only; keep only if ablation adds value"
    ], note="Target: 41.94% → 57-73% RW percentile (to be validated)")
    
    # SLIDE 17: Future Improvements
    img_path = os.path.join(images_dir, "07_future_improvements.png")
    slide = add_image_slide(prs, "Future Improvements & Next Steps", img_path)
    
    # Add expected improvement text with hypothesis qualifier
    exp_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(9), Inches(0.6))
    tf = exp_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Hypothesis / goal range: 41.94% → 57-73% RW percentile (to be validated)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # SLIDE 18: Risks & Mitigations
    add_content_slide(prs, "Risks & Mitigations", [
        "Regime Shifts: BTC market structure changed (institutional adoption)",
        "   Mitigation: Regime-aware model + rolling validation",
        "",
        "Metric Brittleness: SPD sensitive to outliers",
        "   Mitigation: Multiple metrics, sanity gates implemented",
        "",
        "Overfitting: Risk with regime detection on limited data",
        "   Mitigation: Strict train/test split, cross-validation",
        "",
        "Data Quality: Exchange outages, feed gaps",
        "   Mitigation: Polygon as primary, validation checks",
        "",
        "Data Leakage: Causality violations prevented",
        "   Mitigation: Automated tests, embargo periods"
    ])
    
    # SLIDE 19: Deliverables - GRADED FORMAT
    add_two_column_slide(prs, "What We Will Deliver (Graded)",
        left_content=[
            "Reproducible Run:",
            "• Single command execution",
            "• Environment specification",
            "• Deterministic output (seed=42)",
            "",
            "Leakage/Causality Tests:",
            "• Last-row test (PASS)",
            "• Temporal embargo verified",
            "• Zero lookahead confirmed"
        ],
        right_content=[
            "Models & Documentation:",
            "• Midterm: Working baseline + all tests passing",
            "• One improved model (weekly/regime)",
            "• Final report + model card",
            "• Ablation results table",
            "",
            "Code Artifacts:",
            "• Repository with README",
            "• LESSONS_LEARNED.md",
            "• Repro notebook (optional)"
        ]
    )
    
    # Add summary box
    slide = prs.slides[-1]
    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.8))
    tf = summary_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Midterm Result: We falsified a deep-learning approach under strict causality, pivoting to a simpler, more robust baseline"
    p.font.size = Pt(14)
    p.font.color.rgb = GREEN_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # SLIDE 20: The Bottom Line
    add_quote_slide(prs, "The Bottom Line",
        "We built a technically sophisticated system and discovered through rigorous testing that it performs statistically indistinguishable from random. This simplified version represents our evidence-based conclusion: in financial time series, signal quality matters more than model complexity.",
        attribution="Scientific Integrity > Raw Performance")
    
    # SLIDE 21: Questions
    add_content_slide(prs, "Questions & Feedback Request", [
        "We'd appreciate feedback on:",
        "",
        "1. Metric Interpretation: Is our SPD percentile understanding correct?",
        "",
        "2. Split Strategy: Train (2014-2015) / Test (2016-2025) appropriate?",
        "",
        "3. Pivot Scope: Is weekly + regime appropriately scoped?",
        "",
        "4. Baseline Choice: Keep simplified or invest in XGBoost/ensemble?"
    ])
    
    # SLIDE 22: Thank You
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK_BG)
    
    # Thank You title
    ty_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = ty_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Questions subtitle
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(36)
    p.font.color.rgb = YELLOW_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Resources
    res_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    tf = res_box.text_frame
    tf.word_wrap = True
    
    resources = [
        "Documentation: LESSONS_LEARNED.md",
        "Code: Visual_Trading_System repository",
        "Notebook: btc_accumulation_model_simplified.ipynb"
    ]
    for i, res in enumerate(resources):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = res
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    
    # Final quote
    quote_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    tf = quote_box.text_frame
    p = tf.paragraphs[0]
    p.text = '"The best model is the one that works, not the one that sounds impressive."'
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GREEN_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    output_path = "VTS_Midterm_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    main()
