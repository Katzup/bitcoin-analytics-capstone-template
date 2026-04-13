#!/usr/bin/env python3
"""
VTS Midterm Presentation — Revised Edition v2 (post-ChatGPT review)
Visual, concise: 14 slides (was 22). Dark theme.

Slide map:
  1  Title
  2  The Tournament (4-card grid)
  3  Why VTS? (3-column story)
  4  The Pipeline (flowchart)
  5  Reality Check: Scoreboard (table)
  6  Ablation #1 — Allocator Sizing (chart image)
  7  Ablation #2 — The "Aha!" Moment (chart image)
  8  Why GAF+CNN Failed (4-card root cause grid)
  9  What We Learned (Technical) — 3 lesson cards  [NEW]
  10 The Evidence-Based Pivot (side-by-side)
  11 What We Did Right (4 achievement cards)
  12 Roadmap Forward (3-phase timeline)
  13 The Bottom Line (quote + 3 takeaways)
  14 Questions (4 focused Qs)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os, subprocess, sys

# ── Palette ──────────────────────────────────────────────────────────────────
BG     = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy background
CARD   = RGBColor(0x25, 0x25, 0x40)   # card / surface
GREEN  = RGBColor(0x4E, 0xCC, 0xA3)   # #4ecca3 — teal accent
RED    = RGBColor(0xFF, 0x6B, 0x6B)   # #ff6b6b — danger / loss
YELLOW = RGBColor(0xF1, 0xC4, 0x0F)   # #f1c40f — gold / warning
BLUE   = RGBColor(0x3D, 0x9B, 0xE8)   # #3d9be8 — info
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x99, 0x99, 0x99)   # muted text
DKGRN  = RGBColor(0x10, 0x30, 0x25)   # winner highlight bg
DKRED  = RGBColor(0x30, 0x10, 0x10)   # loser highlight bg


# ── Primitive helpers ─────────────────────────────────────────────────────────

def _bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _rect(slide, left, top, w, h, fill=CARD, line_color=None, line_w=0):
    """Plain rectangle."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(w), Inches(h)
    )
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color; s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s


def _card(slide, left, top, w, h, border=GREEN, fill=CARD, bw=1.5):
    """Rounded-rect card — returns shape so caller can populate text_frame."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(w), Inches(h)
    )
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(bw)
    s.text_frame.word_wrap = True
    return s


def _tb(slide, text, left, top, w, h,
        size=14, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT):
    """Textbox helper."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.italic = italic
    p.font.color.rgb = color; p.alignment = align
    return tb


def _hbar(slide, top, color=GREEN, height=0.04):
    """Thin horizontal bar / divider."""
    _rect(slide, 0, top, 10, height, fill=color)


def _section_title(slide, text, top=0.28):
    _tb(slide, text, 0.5, top, 9, 0.72, size=32, bold=True, color=GREEN)
    _hbar(slide, top + 0.74, GREEN, 0.025)


def _card_with_header_body(slide, left, top, w, h,
                            header, body,
                            border=GREEN, fill=CARD,
                            header_size=15, body_size=13,
                            header_color=None, body_color=WHITE):
    """Card with bold header line + body text."""
    if header_color is None:
        header_color = border
    c = _card(slide, left, top, w, h, border=border, fill=fill)
    tf = c.text_frame
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(header_size); p.font.bold = True
    p.font.color.rgb = header_color

    for line in body.split("\n"):
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(body_size)
        p2.font.color.rgb = body_color
        p2.space_before = Pt(3)
    return c


def _pill(slide, text, left, top, w=9.2, h=0.52,
          fill=CARD, border=GREEN, bw=1.5,
          size=12, color=WHITE, bold=False):
    """Small informational pill/badge."""
    c = _card(slide, left, top, w, h, border=border, fill=fill, bw=bw)
    tf = c.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.alignment = PP_ALIGN.CENTER
    return c


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """1 — Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    # Top accent stripe
    _rect(slide, 0, 0, 10, 0.18, fill=GREEN)
    # Bottom stripe
    _rect(slide, 0, 7.32, 10, 0.18, fill=CARD)

    _tb(slide, "When ML Meets Market Reality",
        0.5, 1.9, 9, 1.1,
        size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _tb(slide, "Stacking SATS:  How a 296K-Parameter CNN Got Beat by a Coin Flip",
        0.5, 3.15, 9, 0.85,
        size=22, italic=True, color=YELLOW, align=PP_ALIGN.CENTER)

    _hbar(slide, 4.2, GRAY, 0.02)

    _tb(slide, "GT OMSA Practicum  ·  Trilemma Foundation  ·  Bitcoin Accumulation Tournament",
        0.5, 4.4, 9, 0.45,
        size=14, color=GRAY, align=PP_ALIGN.CENTER)

    _tb(slide, "Robert Katz  ·  rkatz8@gatech.edu",
        0.5, 5.4, 9, 0.45,
        size=17, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    _tb(slide, "Spring 2026",
        0.5, 6.1, 9, 0.4,
        size=13, color=GRAY, align=PP_ALIGN.CENTER)


def slide_02_tournament(prs):
    """2 — The Tournament Challenge (4-card grid)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "The Tournament: Daily Bitcoin Allocation")

    data = [
        (GREEN,  "📊  OBJECTIVE",
         "Predict daily BTC allocation weights\nacross 2016–2025 (~3,300 trading days)"),
        (YELLOW, "🎯  METRIC",
         "Recency-Weighted SPD Percentile\nTarget: top 60% of all submissions"),
        (RED,    "⚖️  CONSTRAINTS",
         "Σ wᵢ = 1.0  ·  wᵢ ≥ 1e-5\nNo shorting  ·  Full budget each day"),
        (BLUE,   "🚫  CRITICAL RULE",
         "Zero lookahead — strict causality\nFeatures use only data through day t"),
    ]

    positions = [(0.35, 1.4), (5.15, 1.4), (0.35, 4.05), (5.15, 4.05)]
    for (left, top), (accent, hdr, body) in zip(positions, data):
        _card_with_header_body(
            slide, left, top, 4.4, 2.35,
            hdr, body,
            border=accent, fill=CARD,
            header_size=15, body_size=14
        )


def slide_03_why_vts(prs):
    """3 — Why Visual Trading Systems? (3-column story)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "Why Visual Trading Systems?")

    cols = [
        (GREEN,  "📚  FINANCE ROOTS",
         "MBA-era chart reading\n(Magee & Edwards)\n\n"
         "Human traders 'see' patterns\nthat quants miss"),
        (YELLOW, "🔬  RESEARCH SIGNAL",
         "Industry reports show\nhigh F1 on candlestick\nclassification tasks\n\n"
         "Explored in OMSA with\nProf. Baich"),
        (BLUE,   "❓  THE HYPOTHESIS",
         "Pattern recognition ≠\nProfitable trading\n\n"
         "Does academic success\nsurvive strict causality?"),
    ]

    for i, (accent, hdr, body) in enumerate(cols):
        _card_with_header_body(
            slide, 0.35 + i * 3.12, 1.4, 3.0, 4.8,
            hdr, body, border=accent,
            header_size=15, body_size=13
        )

    # Hypothesis banner at bottom
    c = _card(slide, 0.35, 6.45, 9.3, 0.72, border=GREEN, fill=DKGRN, bw=1)
    tf = c.text_frame; p = tf.paragraphs[0]
    p.text = ("Hypothesis: GAF images encode price patterns that a CNN can exploit "
              "for superior BTC accumulation under strict causality constraints")
    p.font.size = Pt(13); p.font.italic = True
    p.font.color.rgb = GREEN; p.alignment = PP_ALIGN.CENTER


def slide_04_pipeline(prs):
    """4 — The Visual Trading Pipeline (flowchart)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "Our Initial Approach: The Pipeline")

    steps = [
        ("BTC Price Data",         "90-day rolling window",        GREEN),
        ("GAF Image Generation",   "Polar-coordinate encoding",     GREEN),
        ("Deep CNN (296K params)", "Pattern classification",        YELLOW),
        ("Temperature Calibration","Probability calibration",       YELLOW),
        ("Tilt → EMA → Normalize", "Allocation weights",           GREEN),
    ]

    box_left = 1.0
    box_w    = 4.8

    for i, (name, detail, accent) in enumerate(steps):
        top = 1.35 + i * 1.08
        # Step box
        c = _card(slide, box_left, top, box_w, 0.8, border=accent, fill=CARD, bw=1.5)
        tf = c.text_frame; p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16); p.font.bold = True
        p.font.color.rgb = accent; p.alignment = PP_ALIGN.CENTER

        # Detail label to the right
        _tb(slide, f"  ← {detail}",
            box_left + box_w + 0.2, top + 0.12, 3.6, 0.55,
            size=13, color=GRAY)

        # Down arrow (except last)
        if i < len(steps) - 1:
            _tb(slide, "↓",
                box_left + box_w / 2 - 0.3, top + 0.8, 0.6, 0.32,
                size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # Output pill
    _pill(slide, "OUTPUT: Normalized daily weights  ·  Σw = 1.0  ·  seed = 42  ·  Zero lookahead",
          0.35, 6.85, 9.3, 0.47,
          fill=DKGRN, border=GREEN, size=12, color=GREEN, bold=True)


def slide_05_scoreboard(prs):
    """5 — Reality Check: Scoreboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "Reality Check: First Results")

    _tb(slide,
        "The 296K-parameter CNN barely outperforms a constant neutral signal — "
        "and both fall well short of the 60% target",
        0.5, 1.1, 9, 0.52,
        size=14, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ── Manual table (shapes) ──────────────────────────────────────────────
    headers   = ["Metric",         "CNN (296K)",   "Simplified",   "Target"]
    col_lefts = [0.35,             2.95,           5.35,           7.70]
    col_widths= [2.55,             2.35,           2.30,           1.60]
    row_h     = 0.64
    hdr_top   = 1.85

    # Header row — color tinting
    hdr_fills  = [CARD,  DKRED,    DKGRN,   CARD]
    hdr_colors = [GRAY,  RED,      GREEN,   GRAY]
    for j in range(4):
        c = _rect(slide, col_lefts[j], hdr_top, col_widths[j], 0.50,
                  fill=hdr_fills[j])
        _tb(slide, headers[j],
            col_lefts[j] + 0.05, hdr_top + 0.05,
            col_widths[j] - 0.1, 0.40,
            size=14, bold=True, color=hdr_colors[j],
            align=PP_ALIGN.CENTER)

    rows = [
        ["RW SPD Percentile", "41.43%",   "41.94%",   "> 60%"],
        ["Win Rate vs DCA",   "54.32%",   "70.42%",   "> 50%"],
        ["Execution Time",    "~17 min",  "~2.5 min", "—"],
        ["Model Size",        "1.1 MB",   "None",     "—"],
        ["Reproducibility",   "seed=42",  "seed=42",  "✓"],
    ]
    row_fills = [RGBColor(0x1E, 0x1E, 0x35), CARD]

    for i, row in enumerate(rows):
        top = hdr_top + 0.50 + i * row_h
        for j, cell in enumerate(row):
            fill = DKGRN if j == 2 else row_fills[i % 2]
            _rect(slide, col_lefts[j], top, col_widths[j], row_h - 0.04, fill=fill)
            txt_color = GREEN if j == 2 else (RED if j == 1 else WHITE)
            _tb(slide, cell,
                col_lefts[j] + 0.05, top + 0.10,
                col_widths[j] - 0.1, row_h - 0.15,
                size=17, bold=(j in (1, 2)), color=txt_color,
                align=PP_ALIGN.CENTER)

    _pill(slide,
          "📌  Neutral baseline = constant prob_up = 0.5  →  allocator pipeline unchanged; no model involved",
          0.35, 6.02, 9.3, 0.42,
          fill=RGBColor(0x10, 0x28, 0x20), border=GREEN, size=12,
          color=GREEN, bold=False)

    _pill(slide,
          "⚠️  RW-SPD rewards when and how much you deviate from neutral — not just direction",
          0.35, 6.50, 9.3, 0.42,
          fill=RGBColor(0x28, 0x20, 0x08), border=YELLOW, size=12,
          color=YELLOW, bold=False)


def slide_06_ablation1(prs, images_dir):
    """6 — Ablation #1: Allocator Sizing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "Ablation #1 — Allocator Sensitivity")

    img_path = os.path.join(images_dir, "01_ablation_sensitivity.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(9.0))

    _pill(slide,
          "Verdict: Sweeping allocator aggressiveness moves RW percentile < 0.4 pp  →  NOT the root cause",
          0.35, 6.87, 9.3, 0.47,
          fill=RGBColor(0x28, 0x20, 0x08), border=YELLOW, size=13,
          color=YELLOW, bold=True)


def slide_07_ablation2(prs, images_dir):
    """7 — Ablation #2: The "Aha!" Moment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, 'Ablation #2 — The "Aha!" Moment')

    img_path = os.path.join(images_dir, "02_ablation_signal_quality.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(9.0))

    # Two verdict pills side by side
    _pill(slide,
          "🚨  CNN signal ≈ neutral baseline  →  no predictive edge detected",
          0.35, 6.87, 4.55, 0.47,
          fill=DKRED, border=RED, size=12, color=RED, bold=True)

    _pill(slide,
          "✅  Neutral: +16 pp win rate  ·  simpler  ·  7× faster",
          4.95, 6.87, 4.70, 0.47,
          fill=DKGRN, border=GREEN, size=12, color=GREEN, bold=True)


def slide_08_root_causes(prs, images_dir):
    """8 — Why GAF+CNN Failed (4-card root cause grid, no image)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "Midterm Diagnosis: Why GAF + CNN Failed")

    # Always use the 4-card layout (image removed: contained date inconsistency
    # and "Impact Scores" axis that appeared quantitative but was qualitative).
    causes = [
        (RED,    "⏰  TIME HORIZON MISMATCH",
         "GAF lookback: 90 days\nBTC cycles: 6–18 months\nMisses bull/bear regime shifts"),
        (YELLOW, "📉  DAILY NOISE",
         "Daily prices = low SNR\nTournament forces daily decisions\nAmplifies noise, destroys signal"),
        (BLUE,   "🏗️  WRONG ARCHITECTURE",
         "CNN exploits spatial patterns\nFinance needs temporal modeling\nTranslation invariance irrelevant"),
        (GRAY,   "📅  TRAINING MISMATCH",
         "Train window: 2014–2015 (available historical)\nTest: 2016–2025 (tournament span)\n\nOut-of-distribution: regime shift post-2016"),
    ]

    positions = [(0.35, 1.4), (5.15, 1.4), (0.35, 4.05), (5.15, 4.05)]
    for (left, top), (accent, hdr, body) in zip(positions, causes):
        _card_with_header_body(
            slide, left, top, 4.4, 2.35,
            hdr, body, border=accent,
            header_size=14, body_size=13
        )


def slide_09_technical_lessons(prs):
    """9 — What We Learned (Technical) — 3 horizontal lesson cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "What We Learned (Technical)")

    lessons = [
        (YELLOW, "📊  DAILY SNR IS LOW",
         "Daily price moves are dominated\nby noise, not signal\n\n"
         "Any model exploiting day-to-day\npatterns likely fits noise.\n"
         "Weekly or regime-horizon signals\nhave better SNR."),
        (RED,    "🏗️  WRONG ARCHITECTURE",
         "GAF + CNN exploits spatial\npatterns via 2D convolutions\n\n"
         "BTC accumulation needs temporal\nand regime modeling, not\ntranslation invariance.\n"
         "Problem is sequential, not visual."),
        (BLUE,   "⚙️  ALLOCATOR CAN'T\nRESCUE WEAK SIGNAL",
         "Sweeping tilt & EMA params\ngave < 0.5 pp improvement\n\n"
         "Garbage in → garbage out:\nallocator amplifies signal,\nit doesn't create it.\n"
         "Test signal quality first."),
    ]

    for i, (accent, hdr, body) in enumerate(lessons):
        _card_with_header_body(
            slide, 0.35 + i * 3.12, 1.4, 3.0, 5.3,
            hdr, body, border=accent,
            header_size=14, body_size=13
        )

    _pill(slide,
          "Implication: Test 'does this signal beat random?' before optimizing any downstream component",
          0.35, 6.87, 9.3, 0.47,
          fill=RGBColor(0x10, 0x18, 0x30), border=BLUE, size=12,
          color=BLUE, bold=False)


def slide_10_pivot(prs):
    """10 — The Evidence-Based Pivot (side-by-side)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "The Evidence-Based Pivot")

    # ── Left: CNN (loser) ──────────────────────────────────────────────────
    _card(slide, 0.3, 1.3, 4.2, 5.55, border=RED, fill=DKRED, bw=2)
    _tb(slide, "CNN  (296K params)",
        0.4, 1.42, 4.0, 0.5,
        size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
    _hbar(slide, 1.95, RED, 0.02)

    left_rows = [
        ("RW Percentile",  "41.43%",  RED),
        ("Win Rate",       "54.32%",  RED),
        ("Exec Time",      "~17 min", RED),
        ("Model Size",     "1.1 MB",  GRAY),
    ]
    for i, (metric, val, col) in enumerate(left_rows):
        top = 2.05 + i * 1.08
        _tb(slide, metric, 0.5, top, 3.8, 0.38, size=13, color=GRAY)
        _tb(slide, val,    0.5, top + 0.38, 3.8, 0.55,
            size=24, bold=True, color=col, align=PP_ALIGN.CENTER)

    # ── Center VS ──────────────────────────────────────────────────────────
    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.4), Inches(3.15), Inches(1.2), Inches(1.2)
    )
    oval.fill.solid(); oval.fill.fore_color.rgb = CARD
    oval.line.color.rgb = WHITE; oval.line.width = Pt(2)
    tf = oval.text_frame; p = tf.paragraphs[0]
    p.text = "VS"; p.font.size = Pt(22); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

    # ── Right: Simplified (winner) ─────────────────────────────────────────
    _card(slide, 5.5, 1.3, 4.2, 5.55, border=GREEN, fill=DKGRN, bw=2)
    _tb(slide, "Simplified  (neutral)",
        5.6, 1.42, 4.0, 0.5,
        size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    _hbar(slide, 1.95, GREEN, 0.02)

    right_rows = [
        ("RW Percentile",  "41.94%",  GREEN),
        ("Win Rate",       "70.42%",  GREEN),
        ("Exec Time",      "~2.5 min",GREEN),
        ("Model Size",     "None",    GREEN),
    ]
    for i, (metric, val, col) in enumerate(right_rows):
        top = 2.05 + i * 1.08
        _tb(slide, metric, 5.65, top, 3.8, 0.38, size=13, color=GRAY)
        _tb(slide, val,    5.65, top + 0.38, 3.8, 0.55,
            size=24, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Winner banner
    _pill(slide, "🏆  Simplified wins on ALL dimensions",
          5.5, 6.93, 4.2, 0.45,
          fill=DKGRN, border=GREEN, size=13,
          color=GREEN, bold=True)


def slide_11_quality(prs):
    """11 — What We Did Right (4 achievement cards)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "What We Did Right")

    achievements = [
        (GREEN, "✅  STRICT CAUSALITY",
         "Last-row test  ·  Purge test\nShift test — all PASSED\nZero lookahead confirmed"),
        (GREEN, "✅  DETERMINISTIC",
         "seed = 42\nFully reproducible output\nGrader-safe implementation"),
        (GREEN, "✅  SCIENTIFIC RIGOR",
         "Hypothesis → test → conclude\nSystematic ablation studies\nTransparent negative results"),
        (GREEN, "✅  62/62 TESTS PASSING",
         "Unit + integration + ablation\nBudget accounting tests\nAll results documented"),
    ]

    positions = [(0.35, 1.4), (5.15, 1.4), (0.35, 4.05), (5.15, 4.05)]
    for (left, top), (accent, hdr, body) in zip(positions, achievements):
        _card_with_header_body(
            slide, left, top, 4.4, 2.35,
            hdr, body, border=accent,
            header_size=15, body_size=13
        )


def slide_12_roadmap(prs):
    """12 — Roadmap Forward (3-phase timeline)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "Roadmap Forward")

    phases = [
        (GREEN,  "✅  MIDTERM",   "Status: DONE",
         "Neutral baseline\n62/62 tests passing\nAll ablations complete\nCausality verified"),
        (YELLOW, "🔄  PHASE I",   "→ March 2026",
         "Regime-gated allocator\n(HMM / volatility regime)\nWeekly signal horizon\nXGBoost baseline"),
        (BLUE,   "🎯  PHASE II",  "→ Final Submission",
         "Ensemble integration\nWalk-forward backtest\nFinal report + model card\nFull documentation"),
    ]

    for i, (accent, phase, timing, items) in enumerate(phases):
        left = 0.35 + i * 3.12
        c = _card(slide, left, 1.4, 3.0, 5.3, border=accent, fill=CARD)
        tf = c.text_frame; tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = phase; p.font.size = Pt(16); p.font.bold = True
        p.font.color.rgb = accent

        p2 = tf.add_paragraph()
        p2.text = timing; p2.font.size = Pt(12); p2.font.color.rgb = GRAY
        p2.space_before = Pt(2)

        # Separator
        p3 = tf.add_paragraph()
        p3.text = "─────────────────"; p3.font.size = Pt(9)
        p3.font.color.rgb = GRAY; p3.space_before = Pt(4)

        for item in items.split("\n"):
            pi = tf.add_paragraph()
            pi.text = "  • " + item
            pi.font.size = Pt(13); pi.font.color.rgb = WHITE
            pi.space_before = Pt(4)

    _pill(slide,
          "Goal: attempt to exceed 60%+ RW percentile via regime gating + weekly signals (to be validated in Phase I)",
          0.35, 6.87, 9.3, 0.47,
          fill=RGBColor(0x1A, 0x18, 0x08), border=YELLOW, size=12,
          color=YELLOW, bold=False)


def slide_13_bottom_line(prs):
    """13 — The Bottom Line (quote + 3 takeaways)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    _tb(slide, "THE BOTTOM LINE",
        0.5, 0.3, 9, 0.5,
        size=13, color=GRAY, align=PP_ALIGN.CENTER)

    # Large quote box
    c = _card(slide, 0.7, 0.9, 8.6, 2.9, border=GREEN, fill=DKGRN, bw=2)
    tf = c.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        '"We built a technically sophisticated system and discovered through '
        'rigorous testing that it performs statistically indistinguishable '
        'from random under the RW-SPD tournament metric.\n\n'
        'In financial time series, signal quality matters more than model complexity."'
    )
    p.font.size = Pt(19); p.font.italic = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

    _hbar(slide, 4.1, GRAY, 0.02)

    # Three takeaway cards
    takeaways = [
        (GREEN,  "COMPLEXITY ≠ BETTER",   "296K params matched\nby a constant signal"),
        (YELLOW, "PROCESS IS THE VALUE",  "Scientific integrity\neven with null results"),
        (BLUE,   "SIGNAL FIRST",          "Test 'beats random?'\nbefore optimizing"),
    ]
    for i, (accent, hdr, body) in enumerate(takeaways):
        _card_with_header_body(
            slide, 0.35 + i * 3.12, 4.25, 3.0, 2.2,
            hdr, body, border=accent,
            header_size=14, body_size=13, body_color=WHITE
        )


def slide_14_questions(prs):
    """14 — Questions (4 focused questions)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _section_title(slide, "Questions & Discussion", top=0.28)

    qs = [
        (GREEN,  "Metric Interpretation",
         "Is our understanding of RW-SPD percentile vs win rate correct?"),
        (YELLOW, "Dataset Shift",
         "How would you recommend handling regime shifts / non-stationarity under strict causality?"),
        (BLUE,   "Pivot Scope",
         "Is weekly signals + regime detection appropriately scoped for Phase I?"),
        (RED,    "Baseline Choice",
         "Keep simplified neutral baseline, or invest in XGBoost / ensemble?"),
    ]

    for i, (accent, hdr, body) in enumerate(qs):
        top = 1.35 + i * 1.3
        c = _card(slide, 0.45, top, 9.1, 1.1, border=accent, fill=CARD, bw=1.5)
        tf = c.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Q{i+1}  ·  {hdr}"
        p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = accent
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(13); p2.font.color.rgb = WHITE; p2.space_before = Pt(4)

    _tb(slide,
        "rkatz8@gatech.edu  ·  github.com/Katzup/bitcoin-analytics-capstone-template",
        0.5, 6.92, 9, 0.4,
        size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ── PDF export ────────────────────────────────────────────────────────────────

def export_pdf(pptx_path):
    """Convert PPTX → PDF via LibreOffice (macOS/Linux) or pptx2pdf."""
    abs_path  = os.path.abspath(pptx_path)
    out_dir   = os.path.dirname(abs_path)
    pdf_path  = abs_path.replace(".pptx", ".pdf")

    candidates = [
        "soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/local/bin/soffice",
    ]
    for cmd in candidates:
        try:
            result = subprocess.run(
                [cmd, "--headless", "--convert-to", "pdf",
                 "--outdir", out_dir, abs_path],
                capture_output=True, text=True, timeout=90
            )
            if result.returncode == 0 and os.path.exists(pdf_path):
                print(f"  ✅ PDF: {pdf_path}")
                return pdf_path
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    # Fallback: pptx2pdf (pip install pptx2pdf)
    try:
        from pptx2pdf import convert
        convert(pptx_path, pdf_path)
        print(f"  ✅ PDF: {pdf_path}")
        return pdf_path
    except (ImportError, Exception):
        pass

    print("  ⚠️  PDF export skipped — LibreOffice not found.")
    print(f"      Manual: soffice --headless --convert-to pdf {pptx_path}")
    return None


# ── Main ──────────────────────────────────────────────────────────────────────

def main(out_dir=None):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    images_dir = os.path.join(os.path.dirname(__file__), "images")

    print("Building revised VTS presentation (14 slides)…")
    slide_01_title(prs);                      print("  ✓  1 — Title")
    slide_02_tournament(prs);                 print("  ✓  2 — Tournament")
    slide_03_why_vts(prs);                   print("  ✓  3 — Why VTS?")
    slide_04_pipeline(prs);                  print("  ✓  4 — Pipeline")
    slide_05_scoreboard(prs);                print("  ✓  5 — Scoreboard")
    slide_06_ablation1(prs, images_dir);     print("  ✓  6 — Ablation #1")
    slide_07_ablation2(prs, images_dir);     print("  ✓  7 — Ablation #2")
    slide_08_root_causes(prs, images_dir);   print("  ✓  8 — Root Causes")
    slide_09_technical_lessons(prs);         print("  ✓  9 — Technical Lessons")
    slide_10_pivot(prs);                     print("  ✓ 10 — Pivot")
    slide_11_quality(prs);                   print("  ✓ 11 — Quality")
    slide_12_roadmap(prs);                   print("  ✓ 12 — Roadmap")
    slide_13_bottom_line(prs);               print("  ✓ 13 — Bottom Line")
    slide_14_questions(prs);                 print("  ✓ 14 — Questions")

    if out_dir is None:
        out_dir = os.path.dirname(__file__)

    pptx_path = os.path.join(out_dir, "VTS_Midterm_Presentation_v2.pptx")
    prs.save(pptx_path)
    print(f"\n  ✅ PPTX: {pptx_path}  ({len(prs.slides)} slides)")

    export_pdf(pptx_path)
    return pptx_path


if __name__ == "__main__":
    main()
